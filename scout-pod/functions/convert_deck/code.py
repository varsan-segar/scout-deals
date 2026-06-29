#input_type_name: ConvertDeckInput
#output_type_name: ConvertDeckResult
#function_name: convert_deck

from pydantic import BaseModel
from lemma_sdk import FunctionContext, Pod
import os
import tempfile

class ConvertDeckInput(BaseModel):
    file_path: str

class ConvertDeckResult(BaseModel):
    new_file_path: str
    status: str

def convert_deck(ctx: FunctionContext, data: ConvertDeckInput) -> ConvertDeckResult:
    pod = Pod.from_env()
    
    # Check if it's already a PDF or not a PPT
    if not (data.file_path.endswith(".ppt") or data.file_path.endswith(".pptx")):
        return ConvertDeckResult(new_file_path=data.file_path, status="passthrough")
    
    # It is a PPT. Download it.
    base_name = os.path.basename(data.file_path)
    name_without_ext = os.path.splitext(base_name)[0]
    
    tmpdir = tempfile.gettempdir()
    local_ppt = os.path.join(tmpdir, f"{name_without_ext}.pptx")
    local_pdf = os.path.join(tmpdir, f"{name_without_ext}.pdf")
    
    try:
        from pptx import Presentation
        from fpdf import FPDF
        
        # Download from Lemma pod
        file_bytes = pod.files.download(data.file_path)
        with open(local_ppt, "wb") as f:
            f.write(file_bytes)
            
        # Convert to PDF
        prs = Presentation(local_ppt)
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)
        
        for idx, slide in enumerate(prs.slides):
            if idx > 0:
                pdf.add_page()
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    pdf.multi_cell(0, 10, txt=text)
            
        pdf.output(local_pdf)
        
        # Upload new PDF
        target_name = f"{name_without_ext}.pdf"
        target_local_path = os.path.join(tmpdir, target_name)
        os.rename(local_pdf, target_local_path)
        
        new_path = f"/decks/pdf/{target_name}"
        pod.files.upload(target_local_path, directory_path="/decks/pdf")
        
        if os.path.exists(target_local_path):
            os.remove(target_local_path)
            
        # Delete old PPT from pod
        pod.files.delete(data.file_path)
        
        return ConvertDeckResult(new_file_path=new_path, status="converted")
        
    except Exception as e:
        # If conversion fails, just return original path
        print(f"Error converting deck: {e}")
        return ConvertDeckResult(new_file_path=data.file_path, status="error")
        
    finally:
        if os.path.exists(local_ppt):
            os.remove(local_ppt)
        if os.path.exists(local_pdf):
            os.remove(local_pdf)
